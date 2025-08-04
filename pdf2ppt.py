import sys
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import os


def pdf_to_ppt(pdf_path, ppt_path):
    # 创建一个新的 PPT 演示文稿
    prs = Presentation()
    # 定义幻灯片的布局
    blank_slide_layout = prs.slide_layouts[6]

    try:
        # 提高 dpi 参数以提升图像分辨率
        images = convert_from_path(pdf_path, dpi=500)

        for image in images:
            # 创建新的幻灯片
            slide = prs.slides.add_slide(blank_slide_layout)
            # 临时保存图像
            image_path = 'temp_image.jpg'
            image.save(image_path, 'JPEG')

            # 插入图像到幻灯片中
            left = top = Inches(0)
            pic = slide.shapes.add_picture(image_path, left, top, width=prs.slide_width, height=prs.slide_height)

            # 删除临时图像文件
            os.remove(image_path)

        # 保存 PPT 文件
        prs.save(ppt_path)
        print(f"转换成功，PPT 已保存到 {ppt_path}")
    except Exception as e:
        print(f"转换过程中出现错误: {e}")


if __name__ == "__main__":
    pdf_to_ppt(sys.argv[1], sys.argv[2])
